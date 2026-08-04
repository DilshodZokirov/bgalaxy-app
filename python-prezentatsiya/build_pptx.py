#!/usr/bin/env python3
"""
PBL taqdimot (v7):
- Yoshlarga mos rasmlar
- Boyroq vizual dizayn
- O‘zbekcha atamalar + (inglizcha)
- Oxirida xulosa
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

W, H = Inches(13.333), Inches(7.5)
ASSETS = Path(__file__).resolve().parent / "assets"

BG = RGBColor(0xF7, 0xF9, 0xFC)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0B, 0x12, 0x20)
MUTED = RGBColor(0x33, 0x41, 0x55)
CYAN = RGBColor(0x02, 0x84, 0xC7)
AMBER = RGBColor(0xC2, 0x41, 0x0C)
CORAL = RGBColor(0xBE, 0x12, 0x3C)
GREEN = RGBColor(0x04, 0x78, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xCB, 0xD5, 0xE1)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SOFT = RGBColor(0xE0, 0xF2, 0xFE)
SOFT2 = RGBColor(0xFE, 0xF3, 0xC7)
SOFT3 = RGBColor(0xD1, 0xFA, 0xE5)
SOFT4 = RGBColor(0xFE, 0xE2, 0xE2)
CODE_BG = RGBColor(0xFF, 0xFF, 0xFF)
CODE_FG = RGBColor(0x0B, 0x12, 0x20)
CODE_COMMENT = RGBColor(0x04, 0x78, 0x57)
PROBLEM_BG = RGBColor(0xFE, 0xF2, 0xF2)
SOLUTION_BG = RGBColor(0xEC, 0xFD, 0xF5)


def set_run(run, size=20, bold=False, color=INK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    solid = rPr.find(qn("a:solidFill"))
    if solid is not None:
        rPr.remove(solid)
    solid = rPr.makeelement(qn("a:solidFill"), {})
    srgb = solid.makeelement(qn("a:srgbClr"), {"val": str(color)})
    solid.append(srgb)
    rPr.insert(0, solid)
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", font)


def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def shape(slide, left, top, width, height, fill, *, line=None, rounded=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    return shp


def circle(slide, left, top, size, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def decor(slide, tone="teal"):
    """Boyroq fon: yumshoq doira va burchak aksentlari."""
    if tone == "teal":
        circle(slide, Inches(11.2), Inches(-0.8), Inches(3.2), SOFT)
        circle(slide, Inches(-0.9), Inches(5.4), Inches(2.6), SOFT3)
        shape(slide, 0, 0, Inches(0.16), H, CYAN)
    elif tone == "coral":
        circle(slide, Inches(11.0), Inches(-1.0), Inches(3.4), SOFT4)
        circle(slide, Inches(-1.0), Inches(5.2), Inches(2.8), SOFT2)
        shape(slide, 0, 0, Inches(0.16), H, CORAL)
    elif tone == "green":
        circle(slide, Inches(11.1), Inches(-0.9), Inches(3.3), SOFT3)
        circle(slide, Inches(-0.8), Inches(5.3), Inches(2.5), SOFT)
        shape(slide, 0, 0, Inches(0.16), H, GREEN)
    else:
        shape(slide, 0, 0, Inches(0.16), H, CYAN)


def textbox(slide, left, top, width, height, paragraphs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    first = True
    for item in paragraphs:
        text, kwargs = (item, {}) if isinstance(item, str) else (item[0], item[1])
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = kwargs.get("align", align)
        p.space_after = Pt(kwargs.get("space_after", 4))
        run = p.add_run()
        run.text = text
        set_run(
            run,
            size=kwargs.get("size", 18),
            bold=kwargs.get("bold", False),
            color=kwargs.get("color", INK),
            font=kwargs.get("font", "Calibri"),
        )
    return box


def card(slide, left, top, width, height, title, body, accent=CYAN, fill=CARD):
    shape(slide, left, top, width, height, fill, line=LINE, rounded=True)
    shape(slide, left, top, Inches(0.12), height, accent)
    textbox(
        slide,
        left + Inches(0.28),
        top + Inches(0.18),
        width - Inches(0.42),
        height - Inches(0.3),
        [
            (title, {"size": 15, "bold": True, "color": NAVY, "space_after": 7}),
            (body, {"size": 12, "color": MUTED, "space_after": 0}),
        ],
    )


def code_box(slide, left, top, width, height, title, lines, accent=CYAN):
    shape(slide, left, top, width, height, CODE_BG, line=accent, rounded=True)
    shape(slide, left, top, width, Inches(0.42), accent, rounded=False)
    textbox(
        slide,
        left + Inches(0.2),
        top + Inches(0.05),
        width - Inches(0.3),
        Inches(0.35),
        [(title, {"size": 13, "bold": True, "color": WHITE, "space_after": 0})],
    )
    paras = []
    for line in lines:
        is_comment = line.strip().startswith("#")
        paras.append((
            line if line else " ",
            {
                "size": 13,
                "color": CODE_COMMENT if is_comment else CODE_FG,
                "space_after": 2,
                "font": "Consolas",
            },
        ))
    textbox(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7), paras)


def pill(slide, left, top, width, height, text, bg=CYAN):
    shape(slide, left, top, width, height, bg, rounded=True)
    textbox(
        slide, left, top, width, height,
        [(text, {"size": 11, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 0})],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def heading(slide, eyebrow, title):
    textbox(slide, Inches(0.55), Inches(0.28), Inches(8.5), Inches(1.05), [
        (eyebrow, {"size": 12, "bold": True, "color": CYAN, "space_after": 4}),
        (title, {"size": 23, "bold": True, "color": NAVY, "space_after": 0}),
    ])


def add_pic(slide, path: Path, left, top, width, height=None):
    if not path.exists():
        return None
    if height is None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def new_slide(prs, tone="teal"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG)
    decor(slide, tone)
    return slide


def image_panel_slide(prs, image_name, tone="teal"):
    """Chap matn, o‘ngda rasm paneli."""
    slide = new_slide(prs, tone)
    # o‘ng panel fon
    shape(slide, Inches(8.35), 0, Inches(5.0), H, NAVY)
    add_pic(slide, ASSETS / image_name, Inches(8.35), 0, Inches(5.0), H)
    overlay = ASSETS / "overlay-soft.png"
    if overlay.exists():
        add_pic(slide, overlay, Inches(8.35), 0, Inches(5.0), H)
    return slide


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ---- 1. Hero ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_pic(s, ASSETS / "hero-youth-coding.jpg", 0, 0, W, H)
    add_pic(s, ASSETS / "overlay-dark.png", 0, 0, W, H)
    shape(s, 0, 0, Inches(0.18), H, CORAL)
    pill(s, Inches(0.7), Inches(1.35), Inches(4.6), Inches(0.4), "MUAMMOGA ASOSLANGAN O‘QITISH (PBL)")
    textbox(s, Inches(0.7), Inches(2.0), Inches(9.5), Inches(4.2), [
        ("Python → Django → Django REST", {"size": 34, "bold": True, "color": WHITE, "space_after": 14}),
        ("1) Python nima? Xususiyatlari va boshqa tillarga nisbatan yutuqlari\n2) Muammo: veb kerak → Django\n3) Muammo: telefon + ish stoli + sayt parallel → DRF", {
            "size": 17, "color": RGBColor(0xE2, 0xE8, 0xF0), "space_after": 16,
        }),
        ("Qoida: avval MUAMMO, keyin YECHIM.", {"size": 18, "bold": True, "color": RGBColor(0xFD, 0xBA, 0x74)}),
    ])
    textbox(s, Inches(0.7), Inches(6.7), Inches(10), Inches(0.4), [
        ("v7 · rasmlar + o‘zbekcha atamalar (inglizchasi qavsda)", {"size": 12, "color": RGBColor(0xCB, 0xD5, 0xE1)}),
    ])

    # ---- 2. Python nima? + image ----
    s = image_panel_slide(prs, "python-concept.jpg", "teal")
    heading(s, "PYTHON · SAVOL", "Python nima?")
    shape(s, Inches(0.5), Inches(1.35), Inches(7.5), Inches(1.7), CARD, line=CYAN, rounded=True)
    textbox(s, Inches(0.75), Inches(1.55), Inches(7.1), Inches(1.4), [
        ("Python — umumiy maqsadli (general-purpose), yuqori darajali dasturlash tili.\nKod yoziladi, keyin talqinchi (interpreter) orqali bajariladi. Maqsad: odam o‘qishi oson bo‘lgan dastur yozish.", {
            "size": 14, "color": NAVY, "space_after": 0,
        }),
    ])
    facts = [
        ("Yuqori darajali", "Past darajadagi xotira buyruqlari bilan emas — masala va mantiq bilan ishlaysiz."),
        ("Talqin qilinadi (interpreted)", "Ko‘pincha alohida exe majburiy emas. .py faylni talqinchi ishga tushiradi."),
        ("Ko‘p uslubli (multi-paradigm)", "Oddiy tartib, obyektli (OOP), funksiyaviy uslub — bir tilda birga."),
        ("Ochiq kodli (open source)", "Erkin, katta jamiyat. Windows / macOS / Linux da bir xil fikrlash."),
    ]
    for i, (t, b) in enumerate(facts):
        card(
            s,
            Inches(0.5) + Inches((i % 2) * 3.75),
            Inches(3.3) + Inches((i // 2) * 1.9),
            Inches(3.55),
            Inches(1.75),
            t, b, (CYAN, AMBER, GREEN, CORAL)[i],
        )

    # ---- 3. Features ----
    s = new_slide(prs, "teal")
    # yuqori rasm lenti
    add_pic(s, ASSETS / "python-concept.jpg", Inches(9.4), Inches(0.2), Inches(3.6), Inches(1.05))
    heading(s, "PYTHON · XUSUSIYATLAR", "Ichida nima bor — oson, lekin teran")
    items = [
        ("Sintaksis (syntax)", "Bo‘sh joy bilan blok (indentatsiya) majburiy. Stil bir xil bo‘ladi, o‘qish osonlashadi."),
        ("O‘zgaruvchan tip (dynamic typing)", "Tipni majburan yozmaysiz. Qiymatga qarab aniqlanadi. Katta loyihada tip ko‘rsatmasi (type hint) ixtiyoriy."),
        ("Ichki boylik (batteries included)", "Standart kutubxona boy: fayl, JSON, tarmoq so‘rovi, vaqt, qidiruv andozasi…"),
        ("Keng ekotizim", "Tashqi ombor (PyPI) orqali minglab paket: Django, FastAPI, Pandas, NumPy…"),
        ("Tez ishlab chiqish", "Kam ortiqcha shablon-kod (boilerplate) → prototip tez chiqadi."),
        ("Ko‘chma (portable)", "Bir xil kod turli tizimlarda ishlashi oson. Izolyatsiya muhiti (virtual environment) bilan boshqariladi."),
    ]
    for i, (t, b) in enumerate(items):
        col, row = i % 3, i // 3
        card(
            s,
            Inches(0.45) + Inches(col * 4.2),
            Inches(1.35) + Inches(row * 2.85),
            Inches(4.0),
            Inches(2.7),
            t, b, (CYAN, AMBER, GREEN, CORAL, CYAN, AMBER)[i],
        )

    # ---- 4. Comparison ----
    s = new_slide(prs, "amber" if False else "teal")
    heading(s, "PYTHON · SOLISHTIRMA", "Boshqa tillarga nisbatan yutuqli tomonlari")
    textbox(s, Inches(0.55), Inches(1.25), Inches(12), Inches(0.35), [
        ("Bu “boshqa tillar yomon” degani emas — Python qayerda yutadi, shu yerda aniq aytamiz.", {
            "size": 13, "color": MUTED,
        }),
    ])
    comps = [
        ("Java / C# ga nisbatan", "Kamroq ortiqcha shablon-kod. Tip va klass e’lonlari kamroq. O‘rganish egri chizig‘i yumshoqroq. Veb va sun’iy intellekt prototipida tezroq."),
        ("C / C++ ga nisbatan", "Xotirani qo‘lda boshqarmaysiz. Ko‘rsatkich (pointer) xatolari kam. Sof tezlik pastroq bo‘lishi mumkin — ko‘p vazifa uchun yetarli."),
        ("JavaScript ga nisbatan", "JS asosan brauzer + Node. Python esa sun’iy intellekt, ma’lumot tahlili, avtomatlashtirish, veb orqa qismida kengroq tanlov."),
        ("Go / Rust ga nisbatan", "Go/Rust tizim tezligi uchun kuchli. Python o‘qiluvchanlik, ekotizim va tez ishlab chiqishda yutadi."),
    ]
    for i, (t, b) in enumerate(comps):
        card(
            s,
            Inches(0.45) + Inches((i % 2) * 6.35),
            Inches(1.75) + Inches((i // 2) * 2.55),
            Inches(6.1),
            Inches(2.4),
            t, b, (CYAN, AMBER, GREEN, CORAL)[i],
        )

    # ---- 5. Code compare ----
    s = new_slide(prs, "teal")
    heading(s, "PYTHON · AMALIY SOLISHTIRMA", "Bir xil vazifa: juft sonlarni olish")
    textbox(s, Inches(0.55), Inches(1.25), Inches(12), Inches(0.35), [
        ("Yutuq shu yerda: kam qator, aniq ma’no, o‘qish oson.", {"size": 14, "color": MUTED}),
    ])
    code_box(
        s, Inches(0.45), Inches(1.75), Inches(6.1), Inches(5.1),
        "Java (uzunroq)",
        [
            "List<Integer> nums = Arrays.asList(1,2,3,4,5);",
            "List<Integer> evens = new ArrayList<>();",
            "for (Integer n : nums) {",
            "    if (n % 2 == 0) {",
            "        evens.add(n);",
            "    }",
            "}",
            "System.out.println(evens);",
        ],
        AMBER,
    )
    code_box(
        s, Inches(6.8), Inches(1.75), Inches(5.95), Inches(5.1),
        "Python (qisqa + aniq)",
        [
            "nums = [1, 2, 3, 4, 5]",
            "evens = [n for n in nums if n % 2 == 0]",
            "print(evens)  # [2, 4]",
            "",
            "# Xulosa:",
            "# bir xil natija",
            "# kamroq kod",
            "# o‘qish osonroq",
        ],
        CYAN,
    )

    # ---- 6. Typing ----
    s = new_slide(prs, "teal")
    heading(s, "PYTHON · TIP TIZIMI", "Tipni kim aniqlaydi?")
    card(
        s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(2.1),
        "Ko‘p tillarda",
        "O‘zgaruvchi ochilganda tip yoziladi: butun son (int), matn (String)…\nTekshiruv oldindan bo‘ladi. Aniqlik yuqori, yozish uzunroq.",
        AMBER,
    )
    card(
        s, Inches(6.8), Inches(1.35), Inches(5.95), Inches(2.1),
        "Python da",
        "Tip majburiy yozilmaydi — qiymatga qarab aniqlanadi.\nTez yozish. Katta loyihada tip ko‘rsatmasi (type hint) qo‘shish mumkin.",
        CYAN,
    )
    code_box(
        s, Inches(0.45), Inches(3.65), Inches(6.1), Inches(3.25),
        "Java",
        [
            "int yosh = 18;",
            "String ism = \"Ali\";",
            "double narx = 12.5;",
            "boolean aktiv = true;",
        ],
        AMBER,
    )
    code_box(
        s, Inches(6.8), Inches(3.65), Inches(5.95), Inches(3.25),
        "Python",
        [
            "yosh = 18",
            "ism = \"Ali\"",
            "narx = 12.5",
            "aktiv = True",
            "# tipni Python o‘zi biladi",
        ],
        CYAN,
    )

    # ---- 7. Practical blocks ----
    s = new_slide(prs, "teal")
    heading(s, "PYTHON · AMALIY ASOSLAR", "Ro‘yxat, lug‘at, funksiya — kundalik vositalar")
    code_box(
        s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.5),
        "Ma’lumot tuzilmalari",
        [
            "talaba = {",
            "  \"ism\": \"Ali\",",
            "  \"ball\": 95,",
            "}",
            "print(talaba[\"ism\"])",
            "",
            "fanlar = [\"Python\", \"Django\"]",
            "fanlar.append(\"REST\")",
            "# lug‘at (dict) = kalit→qiymat",
            "# ro‘yxat (list) = tartibli to‘plam",
        ],
        CYAN,
    )
    code_box(
        s, Inches(6.8), Inches(1.35), Inches(5.95), Inches(5.5),
        "Funksiya (function)",
        [
            "def yigindi(a, b):",
            "    return a + b",
            "",
            "print(yigindi(2, 3))  # 5",
            "",
            "def salom(ism=\"mehmon\"):",
            "    return f\"Salom, {ism}\"",
            "",
            "# Funksiya = takrorlanadigan",
            "# mantiqni bitta joyga yig‘ish",
        ],
        GREEN,
    )

    # ---- 8. Domains + image strip ----
    s = new_slide(prs, "teal")
    add_pic(s, ASSETS / "hero-youth-coding.jpg", Inches(9.4), Inches(0.2), Inches(3.6), Inches(1.05))
    heading(s, "PYTHON · QAYERDA ISHLATILADI?", "Bitta til — bir necha yo‘nalish")
    domains = [
        ("Veb orqa qismi (backend)", "Django, Flask, FastAPI.\nSayt va dasturiy interfeys (API)."),
        ("Ma’lumot tahlili (data)", "Pandas, NumPy.\nJadvallar ustida kuchli tahlil."),
        ("Sun’iy intellekt (AI / ML)", "PyTorch, TensorFlow.\nModel o‘qitish va bashorat."),
        ("Avtomatlashtirish", "Fayl, jadval, pochta, saytdan ma’lumot olish (scraping)."),
        ("Ilmiy hisob", "SciPy, Matplotlib.\nGrafik va matematik modellar."),
        ("Skript / DevOps", "Server skriptlari, yordamchi vositalar, tarmoq so‘rovlari."),
    ]
    for i, (t, b) in enumerate(domains):
        col, row = i % 3, i // 3
        card(
            s,
            Inches(0.45) + Inches(col * 4.2),
            Inches(1.35) + Inches(row * 2.85),
            Inches(4.0),
            Inches(2.7),
            t, b, (CYAN, AMBER, GREEN, CORAL, CYAN, AMBER)[i],
        )

    # ---- 9. PROBLEM 1 + web image ----
    s = image_panel_slide(prs, "web-django.jpg", "coral")
    heading(s, "MUAMMO 1", "Python o‘zi — lekin veb kerak")
    shape(s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.4), PROBLEM_BG, line=CORAL, rounded=True)
    textbox(s, Inches(0.8), Inches(1.7), Inches(7.0), Inches(4.9), [
        ("MUAMMO", {"size": 14, "bold": True, "color": CORAL, "space_after": 10}),
        ("Python da skript yozishni o‘rgandik. Lekin real mahsulot ko‘pincha brauzerda: sayt, kabinet, boshqaruv paneli, forma…", {
            "size": 17, "bold": True, "color": NAVY, "space_after": 14,
        }),
        ("Savol: faqat print() bilan to‘liq veb-tizim bo‘ladimi?\nManzil yo‘naltirish (URL routing), ma’lumotlar bazasi, kirish (login), xavfsizlik, sahifa andozasi — bularni noldan yig‘ish og‘ir.", {
            "size": 14, "color": MUTED, "space_after": 14,
        }),
        ("Shuning uchun boshlang‘ich veb uchun dasturiy asos (framework) tanlaymiz.", {
            "size": 15, "bold": True, "color": AMBER,
        }),
    ])

    # ---- 10. Django solution ----
    s = new_slide(prs, "green")
    add_pic(s, ASSETS / "web-django.jpg", Inches(9.4), Inches(0.2), Inches(3.6), Inches(1.05))
    heading(s, "YECHIM 1 · DJANGO", "Boshlang‘ich veb uchun asosiy tavsiya")
    shape(s, Inches(0.45), Inches(1.3), Inches(12.4), Inches(1.35), SOLUTION_BG, line=GREEN, rounded=True)
    textbox(s, Inches(0.75), Inches(1.5), Inches(11.9), Inches(1.05), [
        ("Django — Python dagi to‘liq veb dasturiy asosi (full-stack framework).\nIchida tayyor: kirish tizimi, boshqaruv paneli, ORM, sahifa andozasi, xavfsizlik.", {
            "size": 15, "bold": True, "color": NAVY, "space_after": 0,
        }),
    ])
    reasons = [
        ("Nima uchun Django?", "Boshlang‘ich uchun yo‘l aniq: loyiha/ilova (Project/App), manzil (URL), boshqaruvchi (View), model, andoza (Template)."),
        ("Nima beradi?", "Kirish/ro‘yxat, admin panel, baza bilan ORM, HTML sahifalar, himoya vositalari (CSRF va boshqalar)."),
        ("Qanday o‘rganamiz?", "Avval so‘rov qanday yurishini tushunamiz. Keyin kichik sayt. Keyin dasturiy interfeys (API)."),
    ]
    for i, (t, b) in enumerate(reasons):
        card(s, Inches(0.45) + Inches(i * 4.2), Inches(2.95), Inches(4.0), Inches(3.8), t, b, (GREEN, CYAN, AMBER)[i], SOLUTION_BG)

    # ---- 11. Request flow ----
    s = new_slide(prs, "teal")
    heading(s, "DJANGO · QANDAY ISHLAYDI?", "Bitta so‘rovning (request) ketma-ketligi")
    steps = [
        ("1. So‘rov (request)", "Brauzer manzil ochadi\nyoki forma yuboradi"),
        ("2. Manzil (URL)", "Django to‘g‘ri yo‘lni\ntopadi (urls.py)"),
        ("3. Boshqaruvchi (View)", "Mantiq ishlaydi:\ntekshiruv, hisob"),
        ("4. Model (Model)", "Kerak bo‘lsa bazadan\no‘qiydi / yozadi"),
        ("5. Andoza (Template)", "HTML sahifa yasaydi\n(klassik usul)"),
        ("6. Javob (response)", "Natija brauzerga\nqaytadi"),
    ]
    for i, (t, b) in enumerate(steps):
        col, row = i % 3, i // 3
        card(
            s,
            Inches(0.45) + Inches(col * 4.2),
            Inches(1.35) + Inches(row * 2.85),
            Inches(4.0),
            Inches(2.7),
            t, b, (CYAN, AMBER, GREEN, CORAL, CYAN, AMBER)[i],
        )

    # ---- 12. MTV ----
    s = new_slide(prs, "teal")
    heading(s, "DJANGO · MTV ANDAZASI", "Uchta rol — eslab qolish oson")
    card(s, Inches(0.45), Inches(1.4), Inches(4.0), Inches(5.3),
         "Model (Model)",
         "Ma’lumot modeli.\n\nMasalan: User, Product, Post.\n\nORM orqali SQL kamroq yoziladi — Python klasslari bilan baza boshqariladi.",
         CYAN)
    card(s, Inches(4.65), Inches(1.4), Inches(4.0), Inches(5.3),
         "Andoza (Template)",
         "Ko‘rinish (HTML).\n\nFoydalanuvchi brauzerda shuni ko‘radi.\n\nBoshqaruvchidan kelgan ma’lumot sahifaga joylanadi.",
         AMBER)
    card(s, Inches(8.85), Inches(1.4), Inches(3.95), Inches(5.3),
         "Boshqaruvchi (View)",
         "Mantiq markazi.\n\nSo‘rov keldi → qaror:\nqaysi model? qaysi andoza?\n\nJavobni tayyorlaydi.",
         GREEN)

    # ---- 13. PROBLEM 2 + multi image ----
    s = image_panel_slide(prs, "multi-platform.jpg", "coral")
    heading(s, "MUAMMO 2", "Faqat sayt yetarli emas — real hayot")
    shape(s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.4), PROBLEM_BG, line=CORAL, rounded=True)
    textbox(s, Inches(0.8), Inches(1.7), Inches(7.0), Inches(4.9), [
        ("MUAMMO", {"size": 14, "bold": True, "color": CORAL, "space_after": 10}),
        ("Biz yozgan kod faqat kompyuter brauzerida emas.\nTelefon ilova + ish stoli dasturi (desktop) + veb — parallel ishlashi kerak.", {
            "size": 16, "bold": True, "color": NAVY, "space_after": 12,
        }),
        ("Klassik Django ko‘pincha HTML qaytaradi.\nMobil/ish stoli ilova esa HTML emas — ma’lumot (odatda JSON) so‘raydi.", {
            "size": 14, "color": MUTED, "space_after": 12,
        }),
        ("Xulosa-muammo: faqat andoza (Template) bilan bitta orqa qismni uchala platformaga to‘g‘ri ulash qiyin. Kuchliroq tushuncha kerak: dasturiy interfeys (API).", {
            "size": 14, "bold": True, "color": AMBER,
        }),
    ])

    # ---- 14. Three clients ----
    s = new_slide(prs, "coral")
    heading(s, "MUAMMO 2 · KO‘RINISH", "Uchta mijoz — bitta ma’lumot")
    for i, (t, b) in enumerate([
        ("Telefon (mobile)", "Mobil ilova.\nHTML emas — dasturiy interfeys (API) kerak."),
        ("Ish stoli (desktop)", "Kompyuter dasturi.\nHam ma’lumot so‘raydi."),
        ("Veb (web)", "Brauzer / yagona sahifa ilova (SPA).\nHTML yoki API."),
    ]):
        card(s, Inches(0.45) + Inches(i * 4.2), Inches(1.4), Inches(4.0), Inches(2.5), t, b, CORAL, PROBLEM_BG)
    shape(s, Inches(0.45), Inches(4.2), Inches(12.4), Inches(2.55), CARD, line=AMBER, rounded=True)
    textbox(s, Inches(0.8), Inches(4.45), Inches(11.8), Inches(2.15), [
        ("Savol: mahsulot/foydalanuvchi ma’lumotini uchalasiga qanday beramiz?", {
            "size": 17, "bold": True, "color": NAVY, "space_after": 10,
        }),
        ("Agar faqat HTML bersak — telefon va ish stoli “tushunmaydi”.\nYechim sifatida Django REST Framework (DRF) tavsiya qilinadi.", {
            "size": 15, "color": MUTED,
        }),
    ])

    # ---- 15. DRF ----
    s = new_slide(prs, "green")
    add_pic(s, ASSETS / "multi-platform.jpg", Inches(9.4), Inches(0.2), Inches(3.6), Inches(1.05))
    heading(s, "YECHIM 2 · DJANGO REST (DRF)", "Bir orqa qism — ko‘p platforma")
    shape(s, Inches(0.45), Inches(1.3), Inches(12.4), Inches(1.4), SOLUTION_BG, line=GREEN, rounded=True)
    textbox(s, Inches(0.75), Inches(1.5), Inches(11.9), Inches(1.1), [
        ("DRF — Django ustida dasturiy interfeys (API) qatlami. Ma’lumotni JSON qilib beradi.\nTelefon, ish stoli, veb — parallel ulanadi. Django o‘rnini bosmaydi — uni kengaytiradi.", {
            "size": 15, "bold": True, "color": NAVY, "space_after": 0,
        }),
    ])
    card(s, Inches(0.45), Inches(2.95), Inches(4.0), Inches(3.8),
         "Nima qiladi?",
         "HTML o‘rniga JSON.\nMasalan: /api/products/\n\nSerializator, ViewSet, yo‘naltirgich (router), autentifikatsiya — API uchun asosiy vositalar.",
         GREEN, SOLUTION_BG)
    card(s, Inches(4.65), Inches(2.95), Inches(4.0), Inches(3.8),
         "Nima uchun kerak?",
         "Bir marta yozilgan biznes-mantiq barcha mijozlarga xizmat qiladi.\n\nReal hayotdagi “parallel ishlash” shu yerda yechiladi.",
         CYAN, SOLUTION_BG)
    card(s, Inches(8.85), Inches(2.95), Inches(3.95), Inches(3.8),
         "O‘rganish tartibi",
         "1) Python asoslari\n2) Django (sayt)\n3) DRF (API)\n\nAvval muammo, keyin yechim — shu tartibda.",
         AMBER, SOLUTION_BG)

    # ---- 16. Conclusion with image ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_pic(s, ASSETS / "conclusion-path.jpg", 0, 0, W, H)
    add_pic(s, ASSETS / "overlay-dark.png", 0, 0, W, H)
    shape(s, 0, 0, Inches(0.18), H, GREEN)
    pill(s, Inches(0.7), Inches(0.7), Inches(2.2), Inches(0.38), "XULOSA", GREEN)
    textbox(s, Inches(0.7), Inches(1.3), Inches(11.8), Inches(0.7), [
        ("Bugungi darsdan olib ketiladigan asoslar", {"size": 26, "bold": True, "color": WHITE, "space_after": 0}),
    ])
    points = [
        ("1", "Python", "Umumiy maqsadli til: o‘qilishi oson, tipni o‘zi aniqlaydi, ekotizimi boy."),
        ("2", "Django", "Boshlang‘ich veb uchun to‘liq dasturiy asos (framework) — asosiy tavsiya."),
        ("3", "Muammo", "Faqat sayt yetarli emas: telefon + ish stoli + veb parallel kerak."),
        ("4", "DRF", "Yechim: bitta dasturiy interfeys (API) orqali barcha platformalar ulanadi."),
    ]
    for i, (n, t, b) in enumerate(points):
        y = Inches(2.15) + Inches(i * 1.15)
        circle(s, Inches(0.75), y + Inches(0.15), Inches(0.55), (CYAN, AMBER, CORAL, GREEN)[i])
        textbox(s, Inches(0.75), y + Inches(0.22), Inches(0.55), Inches(0.4), [
            (n, {"size": 14, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 0}),
        ], align=PP_ALIGN.CENTER)
        shape(s, Inches(1.55), y, Inches(10.8), Inches(1.0), RGBColor(0x0F, 0x17, 0x2A), line=RGBColor(0x33, 0x41, 0x55), rounded=True)
        # make panel slightly see-through look via lighter navy border card
        textbox(s, Inches(1.85), y + Inches(0.15), Inches(10.3), Inches(0.75), [
            (t, {"size": 16, "bold": True, "color": WHITE, "space_after": 3}),
            (b, {"size": 13, "color": RGBColor(0xCB, 0xD5, 0xE1), "space_after": 0}),
        ])

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    out_dir = Path(__file__).resolve().parent
    primary = out_dir / "Python-Django-PBL-v7.pptx"
    build(primary)
    for name in (
        "Python-Django-PBL-v6.pptx",
        "Python-Django-PBL-v5.pptx",
        "Python-Dars-01-YANGI.pptx",
        "Python-Django-Yorqin.pptx",
        "Python-Django-Vaaav.pptx",
    ):
        shutil.copyfile(primary, out_dir / name)
