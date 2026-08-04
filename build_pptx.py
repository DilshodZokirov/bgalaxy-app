#!/usr/bin/env python3
"""Build a high-impact widescreen PPTX about Python + Django."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# 16:9
W, H = Inches(13.333), Inches(7.5)

# Palette — OQ / ochiq fon, to'q matn (proyektorda aniq)
BG = RGBColor(0xFF, 0xFF, 0xFF)       # sof oq
BG2 = RGBColor(0xE8, 0xF6, 0xFF)      # ochiq moviy
CARD = RGBColor(0xF4, 0xFB, 0xFF)     # ochiq kartochka
INK = RGBColor(0x11, 0x18, 0x27)      # deyarli qora
MUTED = RGBColor(0x33, 0x41, 0x55)    # to'q kulrang (o'qiladi)
CYAN = RGBColor(0x02, 0x84, 0xC7)     # kuchli moviy
AMBER = RGBColor(0xEA, 0x58, 0x00)    # kuchli apelsin
CORAL = RGBColor(0xE1, 0x1D, 0x48)    # kuchli qizil
GREEN = RGBColor(0x05, 0x96, 0x69)    # kuchli yashil
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0x7D, 0xD3, 0xFC)     # yorqin chegara
NAVY = RGBColor(0x0C, 0x4A, 0x6E)     # sarlavha


def set_run(run, size=20, bold=False, color=INK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # East Asian / latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.insert(0, el)
        el.set("typeface", font)


def add_bg(slide, color=BG):
    """Haqiqiy slide background — shape orqaga surish XML ni buzardi."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color
    return None


def accent_bar(slide, color=CYAN):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def glow_orb(slide, left, top, size, color):
    # Ochiq fonda pastel doira — yordamchi dekor
    orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    orb.fill.solid()
    orb.fill.fore_color.rgb = color
    orb.line.fill.background()


def textbox(slide, left, top, width, height, paragraphs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    first = True
    for item in paragraphs:
        if isinstance(item, str):
            text, kwargs = item, {}
        else:
            text, kwargs = item[0], item[1]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = kwargs.get("align", align)
        p.space_after = Pt(kwargs.get("space_after", 6))
        run = p.add_run()
        run.text = text
        set_run(
            run,
            size=kwargs.get("size", 20),
            bold=kwargs.get("bold", False),
            color=kwargs.get("color", INK),
            font=kwargs.get("font", "Calibri"),
        )
    return box


def card(slide, left, top, width, height, title, body, accent=CYAN):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = CARD
    shp.line.color.rgb = accent
    shp.line.width = Pt(2.25)
    # accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.14), height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    textbox(
        slide,
        left + Inches(0.32),
        top + Inches(0.22),
        width - Inches(0.45),
        height - Inches(0.35),
        [
            (title, {"size": 18, "bold": True, "color": NAVY, "space_after": 8}),
            (body, {"size": 14, "color": MUTED, "space_after": 0}),
        ],
    )


def pill(slide, left, top, width, height, text, bg=CYAN, fg=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    shp.line.fill.background()
    textbox(
        slide,
        left,
        top,
        width,
        height,
        [(text, {"size": 12, "bold": True, "color": fg, "align": PP_ALIGN.CENTER, "space_after": 0})],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def new_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    accent_bar(slide)
    return slide


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ---- 1 Title ----
    s = new_slide(prs)
    glow_orb(s, Inches(10.2), Inches(-0.8), Inches(4), RGBColor(0xBA, 0xE6, 0xFD))
    glow_orb(s, Inches(-1.2), Inches(5.2), Inches(3.2), RGBColor(0xFE, 0xE4, 0xC7))
    pill(s, Inches(0.7), Inches(1.5), Inches(2.4), Inches(0.38), "PYTHON × DJANGO", CYAN, WHITE)
    textbox(
        s,
        Inches(0.7),
        Inches(2.1),
        Inches(11.5),
        Inches(2.2),
        [
            ("Kompyuterlarni boshqaradigan", {"size": 28, "color": MUTED, "space_after": 4}),
            ("ENG KUCHLI BOSHLANG‘ICH TIL", {"size": 44, "bold": True, "color": NAVY, "space_after": 10}),
            ("Qanday ishlaydi? Django nima? Nima uchun millionlab dasturchilar tanlaydi?", {"size": 18, "color": MUTED}),
        ],
    )
    textbox(
        s,
        Inches(0.7),
        Inches(6.5),
        Inches(10),
        Inches(0.5),
        [("Dasturlashni hali bilmaydiganlar uchun — lekin “vaaav” darajasida", {"size": 14, "color": AMBER})],
    )

    # ---- 2 Hook stats ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(1), [
        ("BIRINCHI “VAAAV”", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("Python — shunchaki maktab tili emas", {"size": 32, "bold": True, "color": NAVY}),
    ])
    stats = [
        ("#1", "Ko‘p reytinglarda\neng ommabop tillardan"),
        ("1 qator", "Boshqa tillarda 5–10\nqator bo‘ladigan ish"),
        ("AI+", "Sun’iy intellektning\nasosiy tili"),
        ("Instagram", "Django (Python)\nustida qurilgan"),
    ]
    for i, (a, b) in enumerate(stats):
        x = Inches(0.7) + Inches(i * 3.1)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.9), Inches(3.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE
        textbox(s, x + Inches(0.2), Inches(2.35), Inches(2.5), Inches(3.2), [
            (a, {"size": 28, "bold": True, "color": CYAN if i % 2 == 0 else AMBER, "space_after": 14}),
            (b, {"size": 16, "color": MUTED}),
        ])

    # ---- 3 What is programming ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.4), Inches(11), Inches(1.2), [
        ("ASOS", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("Dasturlash = kompyuterga aniq buyruq berish", {"size": 30, "bold": True, "color": NAVY}),
    ])
    steps = [
        ("Siz g‘oya o‘ylaysiz", "Masalan: son topish o‘yini"),
        ("Python tilida yozasiz", "Kompyuter tushunadigan qoidalar"),
        ("Kompyuter bajaradi", "Sekundda millionlab amal"),
    ]
    for i, (t, b) in enumerate(steps):
        y = Inches(2.0) + Inches(i * 1.5)
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.7), Inches(0.7))
        n.fill.solid()
        n.fill.fore_color.rgb = CYAN if i < 2 else AMBER
        n.line.fill.background()
        textbox(s, Inches(0.8), y + Inches(0.12), Inches(0.7), Inches(0.5), [
            (str(i + 1), {"size": 18, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 0})
        ], align=PP_ALIGN.CENTER)
        textbox(s, Inches(1.8), y, Inches(10), Inches(1.1), [
            (t, {"size": 22, "bold": True, "color": NAVY, "space_after": 4}),
            (b, {"size": 16, "color": MUTED}),
        ])

    # ---- 4 How Python works BIG ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.1), [
        ("PYTHON QANDAY ISHLAYDI?", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("Siz yozgan kod → mashina tushunadigan yo‘l", {"size": 28, "bold": True, "color": NAVY}),
    ])
    pipeline = [
        (".py fayl", "Siz yozgan\nmatn"),
        ("Compiler", "Bytecode\nga aylantiradi"),
        (".pyc", "Tezroq\nyuklanadigan shakl"),
        ("PVM", "Python Virtual\nMachine bajaradi"),
        ("Natija", "Ekran / sayt\n/ hisob"),
    ]
    for i, (t, b) in enumerate(pipeline):
        x = Inches(0.45) + Inches(i * 2.55)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.3), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE
        top_c = CYAN if i != 4 else AMBER
        top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), Inches(2.3), Inches(0.12))
        top.fill.solid()
        top.fill.fore_color.rgb = top_c
        top.line.fill.background()
        textbox(s, x + Inches(0.15), Inches(2.55), Inches(2.0), Inches(2.6), [
            (t, {"size": 18, "bold": True, "color": NAVY, "space_after": 12}),
            (b, {"size": 14, "color": MUTED}),
        ])
        if i < len(pipeline) - 1:
            textbox(s, x + Inches(2.05), Inches(3.4), Inches(0.5), Inches(0.4), [
                ("→", {"size": 22, "bold": True, "color": AMBER, "align": PP_ALIGN.CENTER, "space_after": 0})
            ], align=PP_ALIGN.CENTER)

    # ---- 5 Interpreter vs compiler wow ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.1), [
        ("SIR OCHILADI", {"size": 14, "bold": True, "color": AMBER, "space_after": 4}),
        ("Python — sof interpretator emas", {"size": 30, "bold": True, "color": NAVY}),
    ])
    card(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(4.5),
         "Ko‘pchilik o‘ylaydi",
         "“Python qatorma-qator o‘qiydi.”\n\nBu yarim rost.\n\nAslida CPython avval bytecode yasaydi, keyin Virtual Machine uni uchirib bajaradi.\n\nShuning uchun ham tezroq va aqlliroq ishlaydi.",
         CORAL)
    card(s, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5),
         "Haqiqat (CPython)",
         "1) Parser → AST (daraxt)\n2) Compiler → bytecode\n3) PVM → bajarish\n4) Xotira: havola + GC\n\nNatija: yozish oson, lekin ichida jiddiy “motor” bor.",
         CYAN)

    # ---- 6 Memory wow ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.1), [
        ("XOTIRA SEHRI", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("O‘zgaruvchi — quti emas, ishora", {"size": 30, "bold": True, "color": NAVY}),
    ])
    card(s, Inches(0.7), Inches(1.7), Inches(4), Inches(4.6),
         "Boshqa tillarda",
         "Ko‘pincha qiymat to‘g‘ridan-to‘g‘ri “quti”ga yoziladi.\n\nNusxa olish chalkash bo‘lishi mumkin.",
         CORAL)
    card(s, Inches(4.95), Inches(1.7), Inches(4), Inches(4.6),
         "Python da",
         "x = [1,2,3]\n\nx — bu listga HAVOLA.\ny = x desangiz — yangi list emas, XUDDI SHU obyekt.",
         CYAN)
    card(s, Inches(9.2), Inches(1.7), Inches(3.4), Inches(4.6),
         "Nima uchun muhim?",
         "Tezlik + xotira tejash.\n\nKatta ma’lumotni keraksiz ko‘paytirmaydi.\n\nGC axlatni o‘zi yig‘adi.",
         AMBER)

    # ---- 7 Unique powers ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.1), [
        ("BOSHQA TILLAR QIYIN QILADIGAN NARSALAR", {"size": 14, "bold": True, "color": AMBER, "space_after": 4}),
        ("Pythonning “super kuchlari”", {"size": 30, "bold": True, "color": NAVY}),
    ])
    powers = [
        ("AI / Machine Learning", "TensorFlow, PyTorch, scikit-learn — dunyo standarti deyarli Python."),
        ("1 tilda 5 kasb", "Veb + data + AI + avtomatlashtirish + skript — bitta til."),
        ("Juda tez prototip", "G‘oyadan ishlaydigan dasturga soatlar ichida."),
        ("Fan va hisob", "NASA, CERN, universitetlar — tadqiqotning sevimli tili."),
        ("Avtomatlashtirish", "Excel, fayl, brauzer, bot — kundalik ishlarni “robotlashtirish”."),
        ("O‘qilishi = hujjat", "Kodni o‘qib tushunish oson — jamoa tez ishlaydi."),
    ]
    for i, (t, b) in enumerate(powers):
        col, row = i % 3, i // 3
        card(
            s,
            Inches(0.55) + Inches(col * 4.2),
            Inches(1.65) + Inches(row * 2.55),
            Inches(4.0),
            Inches(2.35),
            t,
            b,
            CYAN if (col + row) % 2 == 0 else AMBER,
        )

    # ---- 8 Framework galaxy ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.3), Inches(12), Inches(1.0), [
        ("FRAMEWORK GALAKTIKASI", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("Pythonning eng mashhur “qurollari”", {"size": 28, "bold": True, "color": NAVY}),
    ])
    fws = [
        ("Django", "To‘liq veb platforma\nInstagram darajasi", AMBER),
        ("FastAPI", "Zamonaviy API\nJuda tez", CYAN),
        ("Flask", "Yengil veb\nKichik loyihalar", RGBColor(0x00, 0xB4, 0xFF)),
        ("PyTorch", "AI / deep learning\nTadqiqot + sanoat", CORAL),
        ("Pandas", "Ma’lumotlar jami\nExcel dan kuchli", GREEN),
        ("Selenium", "Brauzer avtomat\nTest va botlar", RGBColor(0xFF, 0x4F, 0x9A)),
    ]
    for i, (t, b, c) in enumerate(fws):
        col, row = i % 3, i // 3
        x = Inches(0.6) + Inches(col * 4.2)
        y = Inches(1.55) + Inches(row * 2.7)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.0), Inches(2.45))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.0), Inches(0.14))
        strip.fill.solid()
        strip.fill.fore_color.rgb = c
        strip.line.fill.background()
        textbox(s, x + Inches(0.25), y + Inches(0.4), Inches(3.5), Inches(1.8), [
            (t, {"size": 24, "bold": True, "color": NAVY, "space_after": 10}),
            (b, {"size": 15, "color": MUTED}),
        ])

    # ---- 9 Django title ----
    s = new_slide(prs)
    glow_orb(s, Inches(9.5), Inches(-1), Inches(5), RGBColor(0xBB, 0xF7, 0xD0))
    pill(s, Inches(0.7), Inches(2.0), Inches(2.0), Inches(0.38), "DJANGO", AMBER, WHITE)
    textbox(s, Inches(0.7), Inches(2.6), Inches(11), Inches(2.5), [
        ("Sayt qurish uchun", {"size": 26, "color": MUTED, "space_after": 6}),
        ("TAYYOR “SHAHAR INFRASTRUKTURASI”", {"size": 36, "bold": True, "color": NAVY, "space_after": 12}),
        ("Login, baza, admin, xavfsizlik — hammasi ichida. Siz g‘oya va mantiqqa e’tibor qaratasiz.", {"size": 18, "color": MUTED}),
    ])

    # ---- 10 Django how it works ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.3), Inches(12), Inches(1.0), [
        ("DJANGO QANDAY ISHLAYDI?", {"size": 14, "bold": True, "color": AMBER, "space_after": 4}),
        ("Bir so‘rovning sayohati (Request Journey)", {"size": 28, "bold": True, "color": NAVY}),
    ])
    journey = [
        ("1. Brauzer", "Foydalanuvchi\ntugmani bosadi"),
        ("2. URL", "Django yo‘lni\ntopadi"),
        ("3. View", "Mantiq ishlaydi\n(Python kodi)"),
        ("4. Model", "Bazadan\nma’lumot"),
        ("5. Template", "HTML\nyasailadi"),
        ("6. Javob", "Sahifa\nqaytadi"),
    ]
    for i, (t, b) in enumerate(journey):
        col, row = i % 3, i // 3
        card(
            s,
            Inches(0.55) + Inches(col * 4.2),
            Inches(1.55) + Inches(row * 2.7),
            Inches(4.0),
            Inches(2.45),
            t,
            b,
            AMBER if i in (2, 3) else CYAN,
        )

    # ---- 11 MTV ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.1), [
        ("DJANGO ARXITEKTURASI", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("MTV — Model · Template · View", {"size": 30, "bold": True, "color": NAVY}),
    ])
    mtv = [
        ("MODEL", "Ma’lumotlar\n(foydalanuvchi, mahsulot…)\nBazaga bog‘lanadi", CYAN),
        ("VIEW", "Miya\nQanday javob berish?\nQanday hisoblash?", AMBER),
        ("TEMPLATE", "Ko‘rinish\nFoydalanuvchi ko‘radigan\nHTML sahifa", CORAL),
    ]
    for i, (t, b, c) in enumerate(mtv):
        x = Inches(0.7) + Inches(i * 4.15)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(3.9), Inches(4.4))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.9), Inches(3.9), Inches(0.9))
        head.fill.solid()
        head.fill.fore_color.rgb = c
        head.line.fill.background()
        textbox(s, x, Inches(2.1), Inches(3.9), Inches(0.6), [
            (t, {"size": 22, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 0})
        ], align=PP_ALIGN.CENTER)
        textbox(s, x + Inches(0.3), Inches(3.2), Inches(3.3), Inches(2.7), [
            (b, {"size": 18, "color": MUTED}),
        ])

    # ---- 12 Django superpowers ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.0), [
        ("DJANGO SUPERPOWERS", {"size": 14, "bold": True, "color": AMBER, "space_after": 4}),
        ("Boshqa yondashuvlarda qo‘lda yoziladigan narsalar", {"size": 26, "bold": True, "color": NAVY}),
    ])
    supers = [
        ("Admin panel", "1 buyruq bilan boshqaruv paneli — mahsulot/user qo‘shish."),
        ("ORM", "SQL ni deyarli yozmasdan baza bilan ishlash."),
        ("Auth", "Ro‘yxatdan o‘tish, parol, ruxsatlar — tayyor."),
        ("Xavfsizlik", "CSRF, XSS himoyasi — “bateryalar ichida”."),
        ("Migrations", "Baza o‘zgarishlarini versiya qilib saqlaydi."),
        ("Scalable", "Kichik blogdan Instagram darajasigacha yo‘l ochiq."),
    ]
    for i, (t, b) in enumerate(supers):
        col, row = i % 3, i // 3
        card(s, Inches(0.55) + Inches(col * 4.2), Inches(1.6) + Inches(row * 2.55),
             Inches(4.0), Inches(2.35), t, b, AMBER if row == 0 else CYAN)

    # ---- 13 Compare frameworks ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.3), Inches(12), Inches(1.0), [
        ("QAYSI FRAMEWORK?", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("Django · Flask · FastAPI", {"size": 28, "bold": True, "color": NAVY}),
    ])
    rows = [
        ("Django", "Katta sayt, admin, tez start", "To‘liq “kombayn”", AMBER),
        ("Flask", "Kichik API / sayt", "Minimal, erkin", CYAN),
        ("FastAPI", "Zamonaviy API, AI backend", "Juda tez + avto-docs", CORAL),
    ]
    for i, (name, use, vibe, c) in enumerate(rows):
        y = Inches(1.6) + Inches(i * 1.7)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(12), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE
        mark = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(0.16), Inches(1.5))
        mark.fill.solid()
        mark.fill.fore_color.rgb = c
        mark.line.fill.background()
        textbox(s, Inches(1.1), y + Inches(0.3), Inches(2.5), Inches(1.0), [
            (name, {"size": 24, "bold": True, "color": NAVY, "space_after": 0})
        ])
        textbox(s, Inches(4.0), y + Inches(0.25), Inches(4.2), Inches(1.1), [
            ("Qachon?", {"size": 12, "bold": True, "color": MUTED, "space_after": 4}),
            (use, {"size": 16, "color": INK, "space_after": 0}),
        ])
        textbox(s, Inches(8.5), y + Inches(0.25), Inches(3.8), Inches(1.1), [
            ("His?", {"size": 12, "bold": True, "color": MUTED, "space_after": 4}),
            (vibe, {"size": 16, "color": INK, "space_after": 0}),
        ])

    # ---- 14 Real world ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.0), [
        ("HAQIQIY DUNYO", {"size": 14, "bold": True, "color": CYAN, "space_after": 4}),
        ("Python + Django — kim ishlatadi?", {"size": 28, "bold": True, "color": NAVY}),
    ])
    brands = [
        ("Instagram", "Django"),
        ("YouTube", "Python"),
        ("Spotify", "Python"),
        ("NASA", "Python"),
        ("Pinterest", "Django"),
        ("Dropbox", "Python"),
    ]
    for i, (a, b) in enumerate(brands):
        col, row = i % 3, i // 3
        x = Inches(0.7) + Inches(col * 4.15)
        y = Inches(1.7) + Inches(row * 2.5)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.95), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE
        textbox(s, x + Inches(0.3), y + Inches(0.55), Inches(3.3), Inches(1.3), [
            (a, {"size": 26, "bold": True, "color": NAVY, "space_after": 8}),
            (b, {"size": 16, "bold": True, "color": CYAN if b == "Python" else AMBER}),
        ])

    # ---- 15 Path ----
    s = new_slide(prs)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(1.0), [
        ("YO‘L XARITA", {"size": 14, "bold": True, "color": AMBER, "space_after": 4}),
        ("Noldan “vaaav” gacha", {"size": 30, "bold": True, "color": NAVY}),
    ])
    roadmap = [
        ("01", "Python asoslari", "print, if, funksiya"),
        ("02", "Loyiha", "O‘yin / todo"),
        ("03", "Django start", "Birinchi sahifa"),
        ("04", "Django pro", "Login + baza + admin"),
        ("05", "Portfolio", "O‘z saytingiz"),
    ]
    for i, (n, t, b) in enumerate(roadmap):
        x = Inches(0.45) + Inches(i * 2.55)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.4), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD
        box.line.color.rgb = LINE
        textbox(s, x + Inches(0.15), Inches(2.5), Inches(2.1), Inches(3.0), [
            (n, {"size": 20, "bold": True, "color": CYAN if i < 3 else AMBER, "space_after": 12}),
            (t, {"size": 16, "bold": True, "color": NAVY, "space_after": 10}),
            (b, {"size": 13, "color": MUTED}),
        ])

    # ---- 16 Close ----
    s = new_slide(prs)
    glow_orb(s, Inches(-1), Inches(-1), Inches(4), RGBColor(0xFE, 0xCD, 0xD3))
    textbox(s, Inches(0.7), Inches(1.8), Inches(12), Inches(3.5), [
        ("XULOSA", {"size": 14, "bold": True, "color": CYAN, "space_after": 10}),
        ("Python — oson ko‘rinadi.\nIchida esa jiddiy motor.", {"size": 34, "bold": True, "color": NAVY, "space_after": 16}),
        ("Django — shu motor bilan butun sayt shaharini qurish.", {"size": 20, "color": MUTED, "space_after": 20}),
        ("Keyingi qadam: birinchi print(\"Salom\") — keyin Django.", {"size": 18, "bold": True, "color": AMBER}),
    ])

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    out_dir = Path(__file__).resolve().parent
    build(out_dir / "Python-Django-Yorqin.pptx")
    # eski nomga ham nusxa (havolalar uchun)
    import shutil
    shutil.copyfile(out_dir / "Python-Django-Yorqin.pptx", out_dir / "Python-Django-Vaaav.pptx")
